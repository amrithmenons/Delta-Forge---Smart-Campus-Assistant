import os
import re
import uuid
import time
import base64
import traceback
from datetime import datetime

from flask import request, jsonify, current_app
from sqlalchemy.exc import SQLAlchemyError

from . import api
from extensions import db
from models import CourseMaterial, User

# Text extraction imports with availability flags
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    
try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from PIL import Image
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

try:
    from openpyxl import load_workbook
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

# Gemini for subject detection
try:
    import google.generativeai as genai
    GEMINI_API_KEY = os.environ.get('GEMINI_API_KEY', '')
    if GEMINI_API_KEY and len(GEMINI_API_KEY) > 10:
        genai.configure(api_key=GEMINI_API_KEY)
        subject_model = genai.GenerativeModel('gemini-2.5-flash')
        GEMINI_AVAILABLE = True
    else:
        GEMINI_AVAILABLE = False
except:
    GEMINI_AVAILABLE = False

# upload root folder
UPLOAD_ROOT = os.path.join(os.path.dirname(os.path.dirname(__file__)), "uploads")
os.makedirs(UPLOAD_ROOT, exist_ok=True)

# Subject keywords for fallback detection
SUBJECT_KEYWORDS = {
    'Mathematics': ['math', 'algebra', 'calculus', 'geometry', 'trigonometry', 'equation', 'theorem', 'derivative', 'integral'],
    'Physics': ['physics', 'force', 'energy', 'motion', 'velocity', 'acceleration', 'quantum', 'mechanics', 'thermodynamics'],
    'Chemistry': ['chemistry', 'chemical', 'molecule', 'atom', 'reaction', 'compound', 'element', 'periodic', 'organic'],
    'Biology': ['biology', 'cell', 'organism', 'genetics', 'evolution', 'ecosystem', 'species', 'anatomy', 'physiology'],
    'Computer Science': ['programming', 'algorithm', 'data structure', 'software', 'code', 'python', 'java', 'database', 'computer'],
    'History': ['history', 'historical', 'century', 'war', 'civilization', 'empire', 'revolution', 'ancient', 'medieval'],
    'English': ['literature', 'grammar', 'writing', 'essay', 'novel', 'poetry', 'shakespeare', 'language', 'composition'],
    'Economics': ['economics', 'market', 'supply', 'demand', 'gdp', 'inflation', 'trade', 'fiscal', 'monetary'],
}


def decode_data_url(data_url):
    if not data_url:
        return None, None
    m = re.match(r'data:(?P<mime>[^;]+);base64,(?P<b64>.+)', data_url, re.S)
    if m:
        mime = m.group('mime')
        b64 = m.group('b64')
    else:
        mime = None
        b64 = data_url
    try:
        raw = base64.b64decode(b64, validate=True)
        return mime, raw
    except Exception:
        try:
            raw = base64.b64decode(b64)
            return mime, raw
        except Exception:
            return None, None


def ensure_dir(path):
    os.makedirs(path, exist_ok=True)


def detect_subject_with_gemini(text_sample, title):
    """Use Gemini to intelligently detect subject from content"""
    if not GEMINI_AVAILABLE:
        return None
    
    try:
        prompt = f"""Analyze this educational material and determine its subject area.

Title: {title}

Content Sample:
{text_sample[:2000]}

Instructions:
- Identify the PRIMARY academic subject (e.g., Mathematics, Physics, Chemistry, Biology, History, etc.)
- Return ONLY the subject name, nothing else
- If multiple subjects, choose the dominant one
- Use standard academic subject names

Subject:"""

        response = subject_model.generate_content(prompt)
        if response and response.text:
            subject = response.text.strip()
            current_app.logger.info(f"🤖 Gemini detected subject: {subject}")
            return subject
    except Exception as e:
        current_app.logger.warning(f"Gemini subject detection failed: {str(e)}")
    
    return None


def detect_subject_fallback(text_sample, title):
    """Fallback subject detection using keyword matching"""
    # Combine title and text for better detection
    combined = f"{title} {text_sample}".lower()
    
    scores = {}
    for subject, keywords in SUBJECT_KEYWORDS.items():
        score = sum(1 for keyword in keywords if keyword in combined)
        if score > 0:
            scores[subject] = score
    
    if scores:
        detected = max(scores, key=scores.get)
        current_app.logger.info(f"📚 Keyword-detected subject: {detected} (score: {scores[detected]})")
        return detected
    
    return "General Studies"


def detect_subject(text_content, title, user_provided_subject=None):
    """
    Detect subject from content using AI or keyword matching
    Priority: user_provided > gemini > keyword_matching > default
    """
    # If user explicitly provided a subject, use it
    if user_provided_subject and user_provided_subject.strip() and user_provided_subject.lower() not in ['auto', 'detect', 'unknown']:
        return user_provided_subject.strip()
    
    # Get text sample for detection (first 3000 chars)
    text_sample = text_content[:3000] if text_content else ""
    
    # Try Gemini first
    gemini_subject = detect_subject_with_gemini(text_sample, title)
    if gemini_subject:
        return gemini_subject
    
    # Fallback to keyword matching
    return detect_subject_fallback(text_sample, title)


def extract_text_from_pdf(file_path):
    """Extract text from PDF - handles both text-based and scanned PDFs"""
    text = ""
    
    # Method 1: pdfplumber (best for text PDFs)
    if PDFPLUMBER_AVAILABLE:
        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages[:50]):  # Limit to first 50 pages for performance
                    page_text = page.extract_text()
                    if page_text:
                        text += f"\n--- Page {i+1} ---\n{page_text}\n"
            
            if text.strip() and len(text) > 100:
                current_app.logger.info(f"✓ Extracted {len(text)} chars using pdfplumber")
                return text.strip()
        except Exception as e:
            current_app.logger.warning(f"pdfplumber failed: {str(e)}")
    
    # Method 2: PyPDF2 fallback
    if PDF_AVAILABLE and not text.strip():
        try:
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                num_pages = min(len(pdf_reader.pages), 50)  # Limit pages
                for page_num in range(num_pages):
                    page = pdf_reader.pages[page_num]
                    page_text = page.extract_text()
                    if page_text:
                        text += f"\n--- Page {page_num+1} ---\n{page_text}\n"
            
            if text.strip() and len(text) > 100:
                current_app.logger.info(f"✓ Extracted {len(text)} chars using PyPDF2")
                return text.strip()
        except Exception as e:
            current_app.logger.warning(f"PyPDF2 failed: {str(e)}")
    
    return text.strip() if text.strip() else None


def extract_text_from_docx(file_path):
    """Extract text from DOCX"""
    if not DOCX_AVAILABLE:
        return "⚠️ python-docx not available"
    
    try:
        doc = DocxDocument(file_path)
        text = []
        
        for para in doc.paragraphs[:500]:  # Limit paragraphs
            if para.text.strip():
                text.append(para.text)
        
        for table in doc.tables[:20]:  # Limit tables
            table_text = "\n[TABLE]\n"
            for row in table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    table_text += row_text + "\n"
            text.append(table_text)
        
        full_text = '\n\n'.join(text)
        
        if full_text.strip():
            current_app.logger.info(f"✓ Extracted {len(full_text)} chars from DOCX")
            return full_text.strip()
            
    except Exception as e:
        current_app.logger.error(f"DOCX extraction failed: {str(e)}")
        return f"❌ DOCX error: {str(e)}"
    
    return None


def extract_text_from_pptx(file_path):
    """Extract text from PowerPoint"""
    if not PPTX_AVAILABLE:
        return "⚠️ python-pptx not available"
    
    try:
        prs = Presentation(file_path)
        text = []
        
        for slide_num, slide in enumerate(prs.slides[:100], 1):  # Limit slides
            slide_text = f"\n--- Slide {slide_num} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"
            
            if slide_text.strip():
                text.append(slide_text)
        
        full_text = '\n'.join(text)
        
        if full_text.strip():
            current_app.logger.info(f"✓ Extracted {len(full_text)} chars from PPTX")
            return full_text.strip()
            
    except Exception as e:
        current_app.logger.error(f"PPTX extraction failed: {str(e)}")
        return f"❌ PPTX error: {str(e)}"
    
    return None


def extract_text_from_txt(file_path):
    """Extract text from TXT file"""
    encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
    
    for encoding in encodings:
        try:
            with open(file_path, 'r', encoding=encoding) as f:
                text = f.read()
            
            if text.strip():
                current_app.logger.info(f"✓ Extracted {len(text)} chars from TXT ({encoding})")
                return text.strip()
                
        except UnicodeDecodeError:
            continue
        except Exception as e:
            current_app.logger.error(f"TXT extraction failed with {encoding}: {str(e)}")
    
    return "❌ Could not read TXT file"


def truncate_content_intelligently(text, max_chars=10000000):
    """
    Truncate content intelligently while keeping it useful
    MySQL LONGTEXT max: 4GB, but we limit to 10MB for practical reasons
    """
    if not text or len(text) <= max_chars:
        return text
    
    current_app.logger.warning(f"⚠️ Content too large ({len(text)} chars), truncating to {max_chars}")
    
    # Try to truncate at a sentence boundary
    truncated = text[:max_chars]
    
    # Find last sentence ending
    last_period = truncated.rfind('.')
    last_newline = truncated.rfind('\n')
    last_boundary = max(last_period, last_newline)
    
    if last_boundary > max_chars - 1000:  # If boundary is close to limit
        truncated = truncated[:last_boundary + 1]
    
    truncated += f"\n\n[... Content truncated. Original size: {len(text)} characters. Extracted first {len(truncated)} characters for analysis ...]"
    
    return truncated


def extract_content_from_file(file_path, file_type, title):
    """Main extraction function - optimized for large files"""
    current_app.logger.info(f"🔍 Extracting content from: {file_path} (type: {file_type})")
    
    file_type_lower = (file_type or '').lower()
    file_ext = os.path.splitext(file_path)[1].lower()
    
    if not file_type_lower:
        file_type_lower = file_ext.replace('.', '')
    
    extracted_text = None
    
    # PDF files
    if file_type_lower in ['pdf', 'application/pdf'] or file_ext == '.pdf':
        extracted_text = extract_text_from_pdf(file_path)
    
    # DOCX files
    elif file_type_lower in ['docx', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'] or file_ext == '.docx':
        extracted_text = extract_text_from_docx(file_path)
    
    # PowerPoint files
    elif file_type_lower in ['pptx', 'ppt', 'application/vnd.openxmlformats-officedocument.presentationml.presentation'] or file_ext in ['.pptx', '.ppt']:
        extracted_text = extract_text_from_pptx(file_path)
    
    # Text files
    elif file_type_lower in ['txt', 'text/plain'] or file_ext == '.txt':
        extracted_text = extract_text_from_txt(file_path)
    
    # Unsupported format
    else:
        return f"⚠️ Unsupported file format: {file_type_lower or file_ext}"
    
    if extracted_text and len(extracted_text) > 50:
        # Truncate if too large (safety measure)
        extracted_text = truncate_content_intelligently(extracted_text, max_chars=10000000)  # 10MB limit
        
        current_app.logger.info(f"✅ Successfully extracted {len(extracted_text)} characters")
        return extracted_text
    else:
        return f"❌ Could not extract text from '{title}'"


@api.route('/upload-material', methods=['POST'])
def upload_material():
    try:
        # Handle multipart form data (optimized for large files)
        if request.content_type and request.content_type.startswith('multipart/form-data'):
            f = request.files.get('file')
            if not f:
                return jsonify({"error": "no file part in form"}), 400

            student_id = request.form.get('student_id')
            title = request.form.get('title') or f.filename or 'Untitled'
            subject = request.form.get('subject', '').strip()  # Can be empty for auto-detection
            client_file_name = f.filename or 'upload.bin'
            mime = f.mimetype or None
            file_size = 0
            
            # Verify user exists
            user = User.query.get(student_id)
            if user is None:
                return jsonify({"error": "student_id not found"}), 400

            subdir = os.path.join(UPLOAD_ROOT, str(student_id))
            ensure_dir(subdir)

            safe_name = client_file_name.replace("/", "_")[:200]
            unique_name = f"{uuid.uuid4().hex[:8]}_{safe_name}"
            file_path = os.path.join(subdir, unique_name)

            # Stream large file to disk (no memory buffering)
            try:
                with open(file_path, "wb") as wf:
                    chunk_size = 8192  # 8KB chunks
                    while True:
                        chunk = f.stream.read(chunk_size)
                        if not chunk:
                            break
                        wf.write(chunk)
                        file_size += len(chunk)
                
                current_app.logger.info(f"📁 Saved file: {file_path} ({file_size} bytes)")
            except Exception as e:
                current_app.logger.exception("Failed to write file")
                return jsonify({"error": "file_write_failed", "msg": str(e)}), 500

            ext = mime.split('/')[-1] if mime else os.path.splitext(client_file_name)[1].replace('.', '')

        else:
            # JSON upload (for smaller files)
            try:
                data = request.get_json(force=True)
            except Exception as e:
                return jsonify({"error": "invalid json", "msg": str(e)}), 400

            student_id = data.get('student_id')
            title = data.get('title', 'Untitled')
            subject = data.get('subject', '').strip()
            client_file_name = data.get('file_name') or (title.replace(" ", "_")[:60] + ".bin")
            file_type = data.get('file_type') or ''
            file_data = data.get('file_data')

            if not student_id or not file_data:
                return jsonify({"error": "student_id and file_data required"}), 400

            mime, raw_bytes = decode_data_url(file_data)
            if raw_bytes is None:
                return jsonify({"error": "invalid file_data (base64)"}), 400

            file_size = len(raw_bytes)
            ext = (file_type or (mime.split('/')[-1] if mime else '')).lower()

            user = User.query.get(student_id)
            if user is None:
                return jsonify({"error": "student_id not found"}), 400

            subdir = os.path.join(UPLOAD_ROOT, str(student_id))
            ensure_dir(subdir)

            safe_name = client_file_name.replace("/", "_")[:200]
            unique_name = f"{uuid.uuid4().hex[:8]}_{safe_name}"
            file_path = os.path.join(subdir, unique_name)

            try:
                with open(file_path, "wb") as wf:
                    wf.write(raw_bytes)
                current_app.logger.info(f"📁 Saved file: {file_path} ({file_size} bytes)")
            except Exception as e:
                current_app.logger.exception("Failed to write file")
                return jsonify({"error": "file_write_failed", "msg": str(e)}), 500

        # Create database record immediately (before extraction)
        mat = CourseMaterial(
            student_id=str(student_id),
            title=title,
            subject=subject or "Processing...",  # Temporary
            file_type=ext,
            file_name=unique_name,
            file_size=file_size,
            processing_status='processing',
            upload_date=datetime.utcnow(),
            created_at=datetime.utcnow()
        )

        db.session.add(mat)
        db.session.commit()
        
        material_id = mat.id
        current_app.logger.info(f"💾 Created material record: ID={material_id}")

        # EXTRACT TEXT FROM FILE
        try:
            current_app.logger.info(f"🔄 Starting text extraction for: {title}")
            
            extracted_content = extract_content_from_file(file_path, ext, title)
            
            if extracted_content and len(extracted_content) > 50:
                mat.content = extracted_content
                mat.processing_status = 'completed'
                
                # AUTO-DETECT SUBJECT if not provided
                detected_subject = detect_subject(extracted_content, title, subject)
                mat.subject = detected_subject
                
                preview = extracted_content[:300].replace('\n', ' ')
                current_app.logger.info(f"✅ Extraction complete ({len(extracted_content)} chars). Subject: {detected_subject}")
            else:
                mat.content = extracted_content or f"⚠️ Extraction returned empty content"
                mat.processing_status = 'failed'
                mat.subject = subject or "Unknown"
                current_app.logger.warning(f"⚠️ Empty extraction for {title}")
            
            db.session.commit()
            
            return jsonify({
                "id": material_id,
                "chunks_created": 1,
                "status": mat.processing_status,
                "subject": mat.subject,
                "content_length": len(mat.content or ""),
                "content_preview": (mat.content or "")[:200],
                "message": "Upload and extraction completed successfully" if mat.processing_status == 'completed' else "Extraction had issues"
            }), 201

        except Exception as e:
            db.session.rollback()
            current_app.logger.exception("❌ Text extraction failed")
            
            mat.processing_status = 'failed'
            mat.content = f"❌ Extraction error: {str(e)}"
            mat.subject = subject or "Unknown"
            db.session.commit()
            
            return jsonify({
                "id": material_id,
                "error": "extraction_failed",
                "msg": str(e),
                "status": "failed"
            }), 500

    except SQLAlchemyError as e:
        db.session.rollback()
        current_app.logger.exception("DB error in upload_material")
        return jsonify({"error": "db_error", "msg": str(e)}), 500
    except Exception as e:
        current_app.logger.exception("Unhandled error in upload_material")
        return jsonify({"error": "internal_error", "msg": str(e)}), 500


@api.route('/reprocess-material/<material_id>', methods=['POST'])
def reprocess_material(material_id):
    """Re-extract content from an existing material file"""
    try:
        mat = CourseMaterial.query.get(material_id)
        
        if not mat:
            return jsonify({"error": "Material not found"}), 404
        
        file_path = os.path.join(UPLOAD_ROOT, str(mat.student_id), mat.file_name)
        
        if not os.path.exists(file_path):
            return jsonify({"error": "File not found on disk"}), 404
        
        current_app.logger.info(f"🔄 Re-processing material: {mat.title}")
        
        mat.processing_status = 'processing'
        db.session.commit()
        
        extracted_content = extract_content_from_file(file_path, mat.file_type, mat.title)
        
        if extracted_content and len(extracted_content) > 50:
            mat.content = extracted_content
            mat.processing_status = 'completed'
            
            # Re-detect subject
            detected_subject = detect_subject(extracted_content, mat.title, mat.subject)
            mat.subject = detected_subject
            
            current_app.logger.info(f"✅ Re-extraction complete. Subject: {detected_subject}")
        else:
            mat.content = extracted_content or "⚠️ Re-extraction failed"
            mat.processing_status = 'failed'
        
        db.session.commit()
        
        return jsonify({
            "id": material_id,
            "status": mat.processing_status,
            "subject": mat.subject,
            "content_length": len(mat.content or ""),
            "message": "Re-extraction completed"
        }), 200
        
    except Exception as e:
        current_app.logger.exception("❌ Re-processing failed")

        return jsonify({"error": str(e)}), 500
